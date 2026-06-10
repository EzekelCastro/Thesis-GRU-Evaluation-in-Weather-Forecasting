"""
generate_presentation.py
Creates thesis_presentation.pptx  — 16-slide widescreen 16:9 deck.
Run:  python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import datetime

# ── Palette (mirrors tech_brief.docx) ─────────────────────────────────────────
NAVY   = RGBColor(0x1f, 0x38, 0x64)
BLUE   = RGBColor(0x24, 0x71, 0xa3)
DARK   = RGBColor(0x2c, 0x3e, 0x50)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xea, 0xf0, 0xfb)
STEEL  = RGBColor(0xad, 0xb5, 0xbd)
MGRAY  = RGBColor(0xd5, 0xd8, 0xdc)
OFFWH  = RGBColor(0xf4, 0xf6, 0xf7)
DKTXT  = RGBColor(0x2c, 0x3e, 0x50)

C_GRU   = RGBColor(0x1f, 0x77, 0xb4)
C_LSTM  = RGBColor(0xff, 0x7f, 0x0e)
C_RNN   = RGBColor(0x2c, 0xa0, 0x2c)
C_LR    = RGBColor(0xd6, 0x27, 0x28)
C_ARIMA = RGBColor(0x94, 0x67, 0xbd)

STAT_B  = RGBColor(0x24, 0x71, 0xa3)
STAT_G  = RGBColor(0x1a, 0x7a, 0x4a)
STAT_P  = RGBColor(0x7d, 0x3c, 0x98)
STAT_R  = RGBColor(0xc0, 0x39, 0x2b)

# Slide geometry
SW   = Inches(13.33)
SH   = Inches(7.5)
ML   = Inches(0.45)
HDR  = Inches(0.72)
FTR  = Inches(0.34)
CY   = HDR + Inches(0.22)
CW   = SW - ML * 2
CH   = SH - HDR - FTR - Inches(0.22)
FONT = "Calibri"
YEAR = datetime.date.today().year


# ── Primitives ─────────────────────────────────────────────────────────────────

def R(slide, l, t, w, h, fill, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def TB(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def run(p, text, size, bold=False, italic=False, color=DKTXT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return r


def T(slide, l, t, w, h, text, size, bold=False, italic=False,
      color=DKTXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = TB(slide, l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run(p, text, size, bold=bold, italic=italic, color=color)
    return tf


def chip(slide, l, t, w, h, text, fill, size=14, bold=True):
    s = R(slide, l, t, w, h, fill)
    tf = s.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, text, size, bold=bold, color=WHITE)
    return s


def block(slide, l, t, w, h, fill, lines, v_anchor=MSO_ANCHOR.MIDDLE):
    """Colored rect with multiple lines of text inside."""
    s = R(slide, l, t, w, h, fill)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    tf.margin_left   = Inches(0.12)
    tf.margin_right  = Inches(0.12)
    tf.margin_top    = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, (txt, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        if i > 0:
            p.space_before = Pt(4)
        run(p, txt, sz, bold=bld, color=col)
    return s


def hf(slide, title, n):
    """Standard header + footer."""
    R(slide, 0, 0, SW, HDR, NAVY)
    R(slide, 0, 0, Inches(0.14), HDR, BLUE)
    T(slide, ML + Inches(0.1), Inches(0.12), SW - ML * 2, Inches(0.55),
      title, 26, bold=True, color=WHITE)
    R(slide, 0, SH - FTR, SW, FTR, DARK)
    T(slide, ML, SH - FTR + Inches(0.05), SW - ML * 2 - Inches(0.6), Inches(0.28),
      "University of the Cordilleras  ·  B.S. Computer Science  ·  "
      "GRU Evaluation in Weather Forecasting",
      10, color=STEEL)
    T(slide, SW - ML - Inches(0.55), SH - FTR + Inches(0.05),
      Inches(0.5), Inches(0.28), str(n), 10, color=STEEL, align=PP_ALIGN.RIGHT)


def new(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bullet_list(slide, l, t, w, h, items, size=16, color=DKTXT, dot_color=BLUE):
    """Bullet list using a text box."""
    tb = TB(slide, l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = Pt(8)
        rd = p.add_run()
        rd.text = "●  "
        rd.font.size = Pt(size - 1)
        rd.font.bold = True
        rd.font.color.rgb = dot_color
        rd.font.name = FONT
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = FONT


def numbered_list(slide, l, t, w, h, items, size=16, color=DKTXT,
                  num_color=BLUE, num_bg=LIGHT):
    """Numbered circle list."""
    tb = TB(slide, l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = Pt(12)
        rn = p.add_run()
        rn.text = f"  {i + 1}.  "
        rn.font.size = Pt(size)
        rn.font.bold = True
        rn.font.color.rgb = num_color
        rn.font.name = FONT
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = FONT


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Full-bleed title slide."""
    sl = new(prs)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = NAVY

    # Eyebrow
    T(sl, ML, Inches(1.0), CW, Inches(0.4),
      "■  THESIS PRESENTATION  ■", 12, bold=True, color=STEEL,
      align=PP_ALIGN.CENTER)

    # Title
    tb = TB(sl, ML, Inches(1.5), CW, Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "Evaluating the Performance of\nGated Recurrent Units\n"
           "for Multi-Parameter Weather Forecasting",
        40, bold=True, color=WHITE)

    # Divider line
    R(sl, ML + Inches(2), Inches(3.75), CW - Inches(4), Inches(0.05), BLUE)

    # Authors
    T(sl, ML, Inches(3.92), CW, Inches(0.45),
      "Castro, E. M.  ·  Millan, D. J. S.  ·  Mondoñedo, J. E.",
      16, color=RGBColor(0xad, 0xc8, 0xe0), align=PP_ALIGN.CENTER)

    # Institution
    T(sl, ML, Inches(4.44), CW, Inches(0.4),
      "University of the Cordilleras  ·  College of CITCS  ·  B.S. Computer Science",
      14, color=STEEL, align=PP_ALIGN.CENTER)

    # Model colour strip at bottom
    strip_y = Inches(5.55)
    strip_h = Inches(0.35)
    models = [("GRU", C_GRU), ("LSTM", C_LSTM), ("SimpleRNN", C_RNN),
              ("Linear Regression", C_LR), ("ARIMA", C_ARIMA)]
    cw = CW / len(models)
    for i, (name, col) in enumerate(models):
        chip(sl, ML + Inches(i * cw / Inches(1)), strip_y, Inches(cw / Inches(1)),
             strip_h, name, col, size=13)

    # Year badge
    T(sl, ML, Inches(6.05), CW, Inches(0.38),
      f"Baguio City, Benguet, Philippines  ·  {YEAR}",
      12, color=STEEL, align=PP_ALIGN.CENTER)

    # Live URL
    T(sl, ML, Inches(6.5), CW, Inches(0.38),
      "🔗  thesis-gru-evaluation-in-weather-forecasting.streamlit.app",
      12, color=RGBColor(0x5d, 0xad, 0xe8), align=PP_ALIGN.CENTER)


def slide_glance(prs):
    """Slide 2 — Research at a Glance."""
    sl = new(prs)
    hf(sl, "Research at a Glance", 2)

    # 4 stat boxes
    stats = [
        ("5",  "Models\nEvaluated",       STAT_B),
        ("4",  "Weather\nVariables",      STAT_G),
        ("2",  "Weather\nStations",       STAT_P),
        ("5+", "Years of\nTraining Data", STAT_R),
    ]
    bw = Inches(2.9)
    bh = Inches(1.7)
    gap = Inches(0.18)
    total = len(stats) * bw + (len(stats) - 1) * gap
    sx = (SW - total) / 2
    for i, (num, label, col) in enumerate(stats):
        x = sx + i * (bw + gap)
        block(sl, x, CY, bw, bh, col, [
            (num,   34, True,  WHITE),
            (label, 13, True,  RGBColor(0xd6, 0xea, 0xf8)),
        ])

    # Pipeline
    steps = [
        ("DATA",       "Meteostat API",          RGBColor(0x15, 0x43, 0x60)),
        ("PREPROCESS", "Clean · Scale · IQR",    RGBColor(0x1a, 0x52, 0x76)),
        ("MODELS",     "5 ML Algorithms",         NAVY),
        ("EVALUATE",   "MSE · MAE · R² · Acc",   RGBColor(0x4a, 0x23, 0x5a)),
        ("FORECAST",   "7-Day Rolling",           RGBColor(0x14, 0x5a, 0x32)),
    ]
    pipe_y = CY + bh + Inches(0.25)
    pipe_h = Inches(1.5)
    n = len(steps)
    arr_w = Inches(0.32)
    box_w = (CW - (n - 1) * arr_w) / n
    for i, (label, sub, col) in enumerate(steps):
        x = ML + i * (box_w + arr_w)
        block(sl, x, pipe_y, box_w, pipe_h, col, [
            (label, 14, True,  WHITE),
            (sub,   11, False, RGBColor(0xad, 0xc8, 0xe0)),
        ])
        if i < n - 1:
            ax = x + box_w
            s = R(sl, ax, pipe_y, arr_w, pipe_h, MGRAY)
            tf = s.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run(p, "▶", 13, color=RGBColor(0x7f, 0x8c, 0x8d))

    # Caption
    T(sl, ML, pipe_y + pipe_h + Inches(0.18), CW, Inches(0.38),
      "End-to-end pipeline: from raw Meteostat API data to 7-day rolling weather forecasts "
      "across two Philippine weather stations.",
      13, italic=True, color=STEEL, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 3 — Problem Statement & Objectives."""
    sl = new(prs)
    hf(sl, "Problem Statement & Objectives", 3)

    col_l = Inches(7.2)
    col_r = CW - col_l - Inches(0.3)
    col_r_x = ML + col_l + Inches(0.3)

    # Left: problem narrative
    T(sl, ML, CY, col_l, Inches(0.42),
      "RESEARCH BACKGROUND", 12, bold=True, color=BLUE)

    tb = TB(sl, ML, CY + Inches(0.48), col_l, Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run(p, "The Philippines is highly vulnerable to weather extremes — typhoons, monsoon "
           "flooding, and dry spells — making accurate local forecasting critical for "
           "disaster preparedness and agriculture.",
        16, color=DKTXT)

    T(sl, ML, CY + Inches(2.65), col_l, Inches(0.42),
      "RESEARCH GAP", 12, bold=True, color=BLUE)

    tb2 = TB(sl, ML, CY + Inches(3.12), col_l, Inches(1.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run(p2, "Existing comparative studies rarely benchmark multiple deep learning "
            "architectures against classical methods on Philippine station-level data "
            "across all four key meteorological variables simultaneously.",
        16, color=DKTXT)

    # Key question callout
    R(sl, ML, CY + Inches(4.88), col_l, Inches(0.9), LIGHT, line=BLUE, lw=Pt(1.5))
    T(sl, ML + Inches(0.2), CY + Inches(4.98), col_l - Inches(0.4), Inches(0.7),
      "\"Does GRU outperform LSTM, SimpleRNN, Linear Regression, and ARIMA "
      "for multi-variable weather forecasting on Philippine station data?\"",
      14, italic=True, color=NAVY)

    # Right: objectives card
    R(sl, col_r_x, CY, col_r, Inches(5.6), NAVY)
    T(sl, col_r_x + Inches(0.2), CY + Inches(0.18), col_r - Inches(0.3), Inches(0.38),
      "RESEARCH OBJECTIVES", 12, bold=True, color=STEEL, align=PP_ALIGN.LEFT)

    objectives = [
        "Fetch and preprocess real daily weather data\n(Meteostat API, 2020–present)",
        "Train and evaluate GRU, LSTM, SimpleRNN,\nLinear Regression, and ARIMA models",
        "Compare models on four variables:\nprcp, temp, wspd, pres",
        "Benchmark across two stations:\nBaguio City and Manila",
        "Provide a web-based interactive system\nfor result exploration and forecasting",
    ]
    for i, obj in enumerate(objectives):
        oy = CY + Inches(0.72) + i * Inches(1.0)
        block(sl, col_r_x + Inches(0.2), oy, Inches(0.46), Inches(0.46),
              BLUE, [(str(i + 1), 16, True, WHITE)])
        tb3 = TB(sl, col_r_x + Inches(0.78), oy - Inches(0.03),
                 col_r - Inches(0.98), Inches(0.58))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        run(p3, obj, 13, color=WHITE)


def slide_data(prs):
    """Slide 4 — Data Sources & Coverage."""
    sl = new(prs)
    hf(sl, "Data Sources & Station Coverage", 4)

    # Station cards
    stations = [
        ("Baguio City", "Station ID: 98328",
         "Luzon Cordillera region\nElevation: ~1,500 m asl\nCool, wet highland climate",
         RGBColor(0x15, 0x43, 0x60)),
        ("Manila", "Station ID: 98425",
         "National Capital Region\nElevation: ~15 m asl\nTropical lowland climate",
         RGBColor(0x4a, 0x23, 0x5a)),
    ]
    card_w = Inches(5.8)
    card_h = Inches(3.6)
    gap    = Inches(0.5)
    sx     = (SW - 2 * card_w - gap) / 2

    for i, (name, sid, desc, col) in enumerate(stations):
        cx = sx + i * (card_w + gap)
        # Header
        R(sl, cx, CY, card_w, Inches(0.82), col)
        T(sl, cx + Inches(0.2), CY + Inches(0.06), card_w - Inches(0.4), Inches(0.42),
          name, 21, bold=True, color=WHITE)
        T(sl, cx + Inches(0.2), CY + Inches(0.52), card_w - Inches(0.4), Inches(0.28),
          sid, 13, color=STEEL)
        # Body
        R(sl, cx, CY + Inches(0.82), card_w, card_h - Inches(0.82), OFFWH,
          line=MGRAY, lw=Pt(0.5))
        tb = TB(sl, cx + Inches(0.28), CY + Inches(1.02),
                card_w - Inches(0.56), card_h - Inches(1.18))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run(p, desc, 15.5, color=DKTXT)

    # Data period banner
    banner_y = CY + card_h + Inches(0.28)
    R(sl, ML, banner_y, CW, Inches(0.88), NAVY)
    T(sl, ML + Inches(0.3), banner_y + Inches(0.12), CW * 0.55, Inches(0.65),
      "Data Period:  January 1, 2020 — Present  (dynamic, fetched at runtime)",
      16, bold=True, color=WHITE)
    T(sl, ML + CW * 0.6, banner_y + Inches(0.12), CW * 0.38, Inches(0.65),
      "~2,190 daily records per station", 16, color=STEEL, align=PP_ALIGN.RIGHT)

    # Source
    T(sl, ML, banner_y + Inches(1.02), CW, Inches(0.38),
      "Source:  Meteostat Open Weather API  (meteostat.net)  ·  Free, open historical meteorological data",
      13, italic=True, color=STEEL, align=PP_ALIGN.CENTER)


def slide_variables(prs):
    """Slide 5 — Target Variables."""
    sl = new(prs)
    hf(sl, "Target Meteorological Variables", 5)

    vars_ = [
        ("🌧", "Precipitation",    "prcp", "mm / day",
         "Daily total rainfall.\nHighly skewed — many zero-rain days;\nextreme peaks during typhoons.",
         RGBColor(0x15, 0x43, 0x60)),
        ("🌡", "Temperature",      "tavg", "°C",
         "Average daily air temperature.\nWell-structured seasonal signal;\nhighly predictable.",
         RGBColor(0xc0, 0x39, 0x2b)),
        ("💨", "Wind Speed",       "wspd", "km / h",
         "Mean daily wind speed.\nModerate variability;\nelevated during typhoon season.",
         RGBColor(0x1a, 0x7a, 0x4a)),
        ("🔵", "Pressure",         "pres", "hPa",
         "Mean sea-level pressure.\nSteadiest variable;\nstrong anti-correlation with rainfall.",
         RGBColor(0x7d, 0x3c, 0x98)),
    ]

    card_w = Inches(2.9)
    card_h = Inches(4.9)
    gap    = Inches(0.27)
    total  = 4 * card_w + 3 * gap
    sx     = (SW - total) / 2

    for i, (icon, name, code, unit, desc, col) in enumerate(vars_):
        cx = sx + i * (card_w + gap)
        # Coloured top band
        R(sl, cx, CY, card_w, Inches(1.22), col)
        T(sl, cx, CY + Inches(0.05), card_w, Inches(0.6),
          icon, 32, align=PP_ALIGN.CENTER, color=WHITE)
        T(sl, cx, CY + Inches(0.72), card_w, Inches(0.42),
          name, 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # White body
        R(sl, cx, CY + Inches(1.22), card_w, card_h - Inches(1.22),
          OFFWH, line=MGRAY, lw=Pt(0.5))
        # Code pill
        chip(sl, cx + Inches(0.55), CY + Inches(1.4), card_w - Inches(1.1),
             Inches(0.38), code, col, size=13.5)
        T(sl, cx + Inches(0.15), CY + Inches(1.88), card_w - Inches(0.3), Inches(0.34),
          f"Unit: {unit}", 13, color=BLUE, align=PP_ALIGN.CENTER)
        tb = TB(sl, cx + Inches(0.15), CY + Inches(2.3), card_w - Inches(0.3),
                Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run(p, desc, 13, italic=True, color=DKTXT)


def slide_models(prs):
    """Slide 6 — Five Models Compared."""
    sl = new(prs)
    hf(sl, "Five Models Compared", 6)

    models = [
        ("GRU",               C_GRU,   "Gated Recurrent Unit",
         "2-layer Bidirectional · hidden=256\ndropout=0.3 · epochs up to 200",
         "Designed to handle long-range dependencies with gating mechanisms. "
         "Bidirectional context captures both past and future patterns in the sequence."),
        ("LSTM",              C_LSTM,  "Long Short-Term Memory",
         "2-layer Bidirectional · hidden=256\ndropout=0.3 · epochs up to 200",
         "Cell-state architecture explicitly designed to retain long-term information. "
         "Benchmark deep learning comparison for GRU."),
        ("SimpleRNN",         C_RNN,   "Simple Recurrent Neural Network",
         "2-layer Unidirectional · hidden=256\ndropout=0.3 · epochs up to 200",
         "Vanilla RNN baseline. Susceptible to vanishing gradients over long sequences. "
         "Highlights the advantage of gated architectures."),
        ("Linear Regression", C_LR,    "Linear Regression",
         "Flattened sequence · multivariate\nscikit-learn",
         "Classical statistical baseline. Assumes linear relationships. "
         "Tests whether non-linear deep learning adds significant value."),
        ("ARIMA",             C_ARIMA, "Auto-Regressive Integrated Moving Average",
         "auto_arima (AIC) · per-variable\npmdarima · seasonal=False",
         "Traditional univariate time-series model. Industry standard for "
         "statistical forecasting. Fitted independently per variable."),
    ]

    # Chips row
    chip_w = CW / len(models)
    for i, (name, col, *_) in enumerate(models):
        chip(sl, ML + i * chip_w, CY, chip_w, Inches(0.5), name, col, size=15)

    # Model cards
    card_h = Inches(4.6)
    for i, (name, col, full, config, desc) in enumerate(models):
        cx = ML + i * chip_w
        R(sl, cx, CY + Inches(0.52), chip_w, card_h, OFFWH, line=MGRAY, lw=Pt(0.4))
        # Full name
        T(sl, cx + Inches(0.1), CY + Inches(0.62), chip_w - Inches(0.2), Inches(0.28),
          full, 10.5, bold=True, color=col)
        # Config badge
        R(sl, cx + Inches(0.1), CY + Inches(0.96), chip_w - Inches(0.2), Inches(0.82), LIGHT)
        T(sl, cx + Inches(0.15), CY + Inches(1.02), chip_w - Inches(0.3), Inches(0.74),
          config, 10, italic=True, color=NAVY)
        # Description
        tb = TB(sl, cx + Inches(0.1), CY + Inches(1.86),
                chip_w - Inches(0.2), Inches(3.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run(p, desc, 12, color=DKTXT)


def slide_architecture(prs):
    """Slide 7 — Deep Learning Architecture."""
    sl = new(prs)
    hf(sl, "Deep Learning Architecture", 7)

    # Left: diagram
    dw = Inches(7.8)
    layer_cols = [
        ("Input\nSequence", RGBColor(0x27, 0x6e, 0x4e), "seq_len × 4 features"),
        ("BiRNN\nLayer 1",  NAVY,                        "hidden=256 × 2 dirs"),
        ("BiRNN\nLayer 2",  NAVY,                        "dropout=0.3"),
        ("Dense\n128",       RGBColor(0x4a, 0x23, 0x5a), "ReLU activation"),
        ("Output\n4 values", RGBColor(0xc0, 0x39, 0x2b), "prcp·temp·wspd·pres"),
    ]
    arr_w = Inches(0.28)
    box_w = (dw - (len(layer_cols) - 1) * arr_w) / len(layer_cols)
    box_h = Inches(1.72)
    by = CY + Inches(0.45)

    for i, (label, col, sub) in enumerate(layer_cols):
        bx = ML + i * (box_w + arr_w)
        block(sl, bx, by, box_w, box_h, col, [
            (label, 14, True,  WHITE),
            (sub,   10.5,  False, RGBColor(0xad, 0xc8, 0xe0)),
        ])
        if i < len(layer_cols) - 1:
            s = R(sl, bx + box_w, by, arr_w, box_h, MGRAY)
            tf = s.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run(p, "▶", 12, color=DKTXT)

    # Shared config label
    T(sl, ML, by + box_h + Inches(0.15), dw, Inches(0.34),
      "GRU & LSTM: Bidirectional  ·  SimpleRNN: Unidirectional  ·  "
      "Early stopping patience=30  ·  Adam optimizer + ReduceLROnPlateau",
      11.5, italic=True, color=STEEL, align=PP_ALIGN.CENTER)

    # Right: parameter table
    tx = ML + dw + Inches(0.3)
    tw = CW - dw - Inches(0.3)
    T(sl, tx, CY, tw, Inches(0.38), "HYPERPARAMETERS", 12, bold=True, color=BLUE)
    params = [
        ("Hidden Size",    "256"),
        ("Layers",         "2 (optimal)"),
        ("Bidirectional",  "GRU, LSTM"),
        ("Dropout",        "0.3"),
        ("Seq. Length",    "45 days"),
        ("Batch Size",     "16"),
        ("Max Epochs",     "200"),
        ("Patience",       "30"),
        ("Optimizer",      "Adam  (lr=1e-3)"),
        ("LR Scheduler",   "ReduceLROnPlateau"),
        ("Loss",           "MSE"),
    ]
    for i, (k, v) in enumerate(params):
        py = CY + Inches(0.42) + i * Inches(0.46)
        bg_c = LIGHT if i % 2 == 0 else WHITE
        R(sl, tx, py, tw, Inches(0.44), bg_c, line=MGRAY, lw=Pt(0.3))
        T(sl, tx + Inches(0.1), py + Inches(0.06), tw * 0.52, Inches(0.35),
          k, 12, bold=True, color=NAVY)
        T(sl, tx + tw * 0.55, py + Inches(0.06), tw * 0.42, Inches(0.35),
          v, 12, color=DKTXT)


def slide_preprocessing(prs):
    """Slide 8 — Data Preprocessing Pipeline."""
    sl = new(prs)
    hf(sl, "Data Preprocessing Pipeline", 8)

    steps = [
        ("1", "Missing Value\nImputation",
         "Forward-fill → Backward-fill\n→ Linear interpolation",
         RGBColor(0x15, 0x43, 0x60)),
        ("2", "Outlier Removal\n(IQR)",
         "Replaces values outside\nQ1 − 1.5·IQR / Q3 + 1.5·IQR\n"
         "Skipped for prcp\n(right-skewed distribution)",
         RGBColor(0x1a, 0x52, 0x76)),
        ("3", "Feature Scaling\n(MinMax)",
         "Per-column MinMaxScaler\nfitted on training set only\n"
         "Scalers saved for\ninverse-transform",
         NAVY),
        ("4", "Train / Test\nSplit",
         "80% train · 20% test\nTemporal order preserved\nNo shuffling",
         RGBColor(0x4a, 0x23, 0x5a)),
        ("5", "Sliding Window\nSequences",
         "seq_len = 45 days\nX: (N, 45, 4)\ny: (N, 4)\nOne step ahead",
         RGBColor(0x14, 0x5a, 0x32)),
    ]

    box_w = Inches(2.3)
    box_h = Inches(4.5)
    arr_w = Inches(0.26)
    total = len(steps) * box_w + (len(steps) - 1) * arr_w
    sx = (SW - total) / 2
    by = CY + Inches(0.12)

    for i, (num, title, desc, col) in enumerate(steps):
        bx = sx + i * (box_w + arr_w)
        # Number badge
        block(sl, bx, by, box_w, Inches(0.46), col,
              [(f"STEP {num}", 12, True, WHITE)])
        # Title
        R(sl, bx, by + Inches(0.46), box_w, Inches(0.95), col)
        T(sl, bx + Inches(0.1), by + Inches(0.55), box_w - Inches(0.2), Inches(0.78),
          title, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Body
        R(sl, bx, by + Inches(1.41), box_w, box_h - Inches(1.41),
          OFFWH, line=MGRAY, lw=Pt(0.4))
        tb = TB(sl, bx + Inches(0.1), by + Inches(1.52),
                box_w - Inches(0.2), box_h - Inches(1.62))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run(p, desc, 12.5, color=DKTXT)

        if i < len(steps) - 1:
            ax = bx + box_w
            s = R(sl, ax, by + box_h / 2 - Inches(0.22), arr_w, Inches(0.44), MGRAY)
            tf2 = s.text_frame
            tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            run(p2, "▶", 12, color=DKTXT)


def slide_ui_overview(prs):
    """Slide 9 — System Features (UI tabs overview)."""
    sl = new(prs)
    hf(sl, "System Features — Interactive Web Application", 9)

    tabs = [
        ("📈  Predictions",
         "Tabs: Predictions",
         [
             "Side-by-side line charts: Actual vs. Predicted",
             "Separate chart per model per variable",
             "Both stations in one run",
             "Colour-coded model legends",
         ], C_GRU),
        ("📊  Metrics",
         "Tabs: Metrics",
         [
             "Numeric table: MSE, MAE, R², Accuracy (%)",
             "Interactive bar charts — switch between metrics",
             "'Best Model' summary table per variable",
             "Colour-coded model bars",
         ], C_LSTM),
        ("🔮  Forecast",
         "Tabs: Forecast",
         [
             "7-day rolling day-by-day forecast",
             "Check-box model & variable selection",
             "Day-labeled table + multi-model line chart",
             "Starts day after last available data",
         ], C_RNN),
        ("⬇️  Download",
         "Tabs: Download",
         [
             "Export full results as Excel workbook",
             "Separate sheet per station",
             "Predictions & metrics included",
             "One-click download button",
         ], C_ARIMA),
    ]

    card_w = Inches(2.95)
    card_h = Inches(5.2)
    gap    = Inches(0.18)
    sx     = (SW - 4 * card_w - 3 * gap) / 2

    for i, (tab_label, sub, bullets, col) in enumerate(tabs):
        cx = sx + i * (card_w + gap)
        # Tab-style header
        R(sl, cx, CY, card_w, Inches(0.6), col)
        T(sl, cx + Inches(0.1), CY + Inches(0.1), card_w - Inches(0.2), Inches(0.44),
          tab_label, 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Body
        R(sl, cx, CY + Inches(0.6), card_w, card_h - Inches(0.6),
          OFFWH, line=MGRAY, lw=Pt(0.5))
        # Bullets
        tb = TB(sl, cx + Inches(0.15), CY + Inches(0.78),
                card_w - Inches(0.3), card_h - Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, bul in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(0 if j == 0 else 10)
            rd = p.add_run()
            rd.text = "▸  "
            rd.font.size = Pt(12.5)
            rd.font.bold = True
            rd.font.color.rgb = col
            rd.font.name = FONT
            rt = p.add_run()
            rt.text = bul
            rt.font.size = Pt(13.5)
            rt.font.color.rgb = DKTXT
            rt.font.name = FONT


def slide_metrics(prs):
    """Slide 10 — Evaluation Metrics."""
    sl = new(prs)
    hf(sl, "Evaluation Metrics", 10)

    metrics = [
        ("Accuracy (%)",
         "Acc = max(0,  100 − MAPE)",
         "MAPE = mean |y − ŷ| / |y|  × 100",
         "Primary interpretability metric. Intuitive percentage score. "
         "Zero-actual values excluded from MAPE to avoid division by zero.",
         STAT_B),
        ("MSE",
         "MSE = (1/n) Σ (y − ŷ)²",
         "Mean Squared Error",
         "Penalises large errors heavily due to squaring. "
         "Sensitive to outlier predictions — relevant for extreme precipitation events.",
         STAT_G),
        ("MAE",
         "MAE = (1/n) Σ |y − ŷ|",
         "Mean Absolute Error",
         "Robust to outliers. Directly interpretable in the original unit "
         "(mm, °C, km/h, hPa). Paired with MSE to diagnose outlier sensitivity.",
         STAT_P),
        ("R²",
         "R² = 1 − Σ(y−ŷ)² / Σ(y−ȳ)²",
         "Coefficient of Determination",
         "Proportion of variance explained. R²=1 is perfect; R²=0 means the model "
         "is no better than predicting the mean; negative values indicate poor fit.",
         STAT_R),
    ]

    card_w = Inches(2.9)
    card_h = Inches(5.1)
    gap    = Inches(0.22)
    sx     = (SW - 4 * card_w - 3 * gap) / 2

    for i, (name, formula, label, desc, col) in enumerate(metrics):
        cx = sx + i * (card_w + gap)
        # Header
        R(sl, cx, CY, card_w, Inches(0.62), col)
        T(sl, cx, CY + Inches(0.1), card_w, Inches(0.44),
          name, 17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Formula box
        R(sl, cx, CY + Inches(0.62), card_w, Inches(0.96), LIGHT)
        T(sl, cx + Inches(0.1), CY + Inches(0.7), card_w - Inches(0.2), Inches(0.58),
          formula, 14, bold=True, color=col, align=PP_ALIGN.CENTER)
        T(sl, cx + Inches(0.1), CY + Inches(1.18), card_w - Inches(0.2), Inches(0.3),
          label, 11, italic=True, color=STEEL, align=PP_ALIGN.CENTER)
        # Description
        R(sl, cx, CY + Inches(1.58), card_w, card_h - Inches(1.58),
          OFFWH, line=MGRAY, lw=Pt(0.4))
        tb = TB(sl, cx + Inches(0.15), CY + Inches(1.74),
                card_w - Inches(0.3), card_h - Inches(1.82))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run(p, desc, 13, color=DKTXT)


def slide_results_predictions(prs):
    """Slide 11 — Predictions Tab Results."""
    sl = new(prs)
    hf(sl, "Results — Predictions Tab", 11)

    # Left description
    lw = Inches(4.2)
    T(sl, ML, CY, lw, Inches(0.42),
      "WHAT THE TAB SHOWS", 12, bold=True, color=BLUE)

    items = [
        "Line chart for each variable per station",
        "Actual values (ground truth from test set)",
        "Predicted values overlaid per model",
        "Models colour-coded consistently across all charts",
        "X-axis: date · Y-axis: physical unit (mm, °C, km/h, hPa)",
        "Separate row per variable, enabling visual diagnosis of\nwhere each model fails or succeeds",
    ]
    bullet_list(sl, ML, CY + Inches(0.5), lw, Inches(3.5), items,
                size=14, color=DKTXT, dot_color=BLUE)

    T(sl, ML, CY + Inches(4.2), lw, Inches(0.42),
      "EXPECTED FINDINGS", 12, bold=True, color=BLUE)
    findings = [
        "Temperature & Pressure: highest accuracy across all models",
        "Precipitation: hardest to predict — sparsity & extremes",
        "GRU / LSTM: outperform RNN and classical methods on most variables",
        "ARIMA: competitive on smooth variables (temp, pres)",
    ]
    bullet_list(sl, ML, CY + Inches(4.68), lw, Inches(1.85), findings,
                size=13, color=DKTXT, dot_color=NAVY)

    # Right: mock chart panel
    pw = CW - lw - Inches(0.35)
    px = ML + lw + Inches(0.35)
    R(sl, px, CY, pw, Inches(5.6), OFFWH, line=MGRAY, lw=Pt(0.5))
    T(sl, px, CY, pw, Inches(0.38),
      "  Sample — Precipitation (prcp)", 13, bold=True, color=NAVY)

    # Simulated chart lines
    chart_y = CY + Inches(0.42)
    chart_h = Inches(2.45)
    R(sl, px + Inches(0.1), chart_y, pw - Inches(0.2), chart_h,
      WHITE, line=MGRAY, lw=Pt(0.3))
    T(sl, px + Inches(0.18), chart_y + Inches(0.12), pw - Inches(0.35), Inches(0.28),
      "Precipitation (mm)", 10, italic=True, color=STEEL)

    # Fake legend
    legend_y = chart_y + chart_h + Inches(0.08)
    legend_items = [
        ("Actual", DARK), ("GRU", C_GRU), ("LSTM", C_LSTM),
        ("RNN", C_RNN), ("LR", C_LR), ("ARIMA", C_ARIMA),
    ]
    lx = px + Inches(0.15)
    for j, (lname, lcol) in enumerate(legend_items):
        R(sl, lx, legend_y + Inches(0.08), Inches(0.24), Inches(0.12), lcol)
        T(sl, lx + Inches(0.28), legend_y, Inches(0.68), Inches(0.3),
          lname, 10, color=DKTXT)
        lx += Inches(0.96)

    # Temperature chart below (smaller)
    T(sl, px, legend_y + Inches(0.34), pw, Inches(0.35),
      "  Sample — Temperature (°C)", 13, bold=True, color=NAVY)
    R(sl, px + Inches(0.1), legend_y + Inches(0.72), pw - Inches(0.2), Inches(1.88),
      WHITE, line=MGRAY, lw=Pt(0.3))
    T(sl, px + Inches(0.18), legend_y + Inches(0.82), pw - Inches(0.35), Inches(0.28),
      "Temperature (°C)", 10, italic=True, color=STEEL)

    T(sl, px + Inches(0.1), legend_y + Inches(2.7), pw - Inches(0.2), Inches(0.34),
      "[ Actual charts generated after running the analysis on the Streamlit app ]",
      11, italic=True, color=STEEL, align=PP_ALIGN.CENTER)


def slide_results_metrics(prs):
    """Slide 12 — Metrics Tab Results."""
    sl = new(prs)
    hf(sl, "Results — Metrics Tab", 12)

    # Left
    lw = Inches(4.2)
    T(sl, ML, CY, lw, Inches(0.42),
      "METRICS TAB FEATURES", 12, bold=True, color=BLUE)

    features = [
        "Numeric table: all 4 metrics × 5 models per variable",
        "Best model highlighted per metric",
        "Radio selector: switch chart between Accuracy / R² / MSE / MAE",
        "Vertical bar chart — one bar per model, colour-coded",
        "X-axis: model names · Y-axis: metric value",
        "Separate chart panel per station",
    ]
    bullet_list(sl, ML, CY + Inches(0.5), lw, Inches(2.85), features,
                size=14, color=DKTXT, dot_color=BLUE)

    T(sl, ML, CY + Inches(3.5), lw, Inches(0.42),
      "INTERPRETATION GUIDE", 12, bold=True, color=BLUE)
    T(sl, ML, CY + Inches(3.98), lw, Inches(1.8),
      "Higher Accuracy (%) and R² → better fit\n"
      "Lower MSE and MAE → smaller prediction error\n\n"
      "Compare across models within the same variable "
      "to identify the strongest performer for each forecasting task.",
      14, italic=True, color=DKTXT)

    # Right: bar chart mockup
    pw = CW - lw - Inches(0.35)
    px = ML + lw + Inches(0.35)
    R(sl, px, CY, pw, Inches(5.6), OFFWH, line=MGRAY, lw=Pt(0.5))
    T(sl, px + Inches(0.12), CY + Inches(0.1), pw - Inches(0.24), Inches(0.34),
      "Accuracy (%)  —  Temperature  |  Baguio City",
      12, bold=True, color=NAVY)

    # Simulated bars
    bar_y_bottom = CY + Inches(5.15)
    bar_max_h    = Inches(3.5)
    bar_w        = Inches(1.0)
    bar_gap      = Inches(0.28)
    bar_vals     = [92, 91, 87, 83, 89]
    bar_colors   = [C_GRU, C_LSTM, C_RNN, C_LR, C_ARIMA]
    bar_names    = ["GRU", "LSTM", "RNN", "LR", "ARIMA"]
    bx_start     = px + Inches(0.3)

    for j, (val, col, name) in enumerate(zip(bar_vals, bar_colors, bar_names)):
        bh = bar_max_h * val / 100
        bx = bx_start + j * (bar_w + bar_gap)
        R(sl, bx, bar_y_bottom - bh, bar_w, bh, col)
        T(sl, bx, bar_y_bottom - bh - Inches(0.34), bar_w, Inches(0.32),
          f"{val}%", 12, bold=True, color=col, align=PP_ALIGN.CENTER)
        T(sl, bx, bar_y_bottom + Inches(0.05), bar_w, Inches(0.32),
          name, 12, bold=True, color=DKTXT, align=PP_ALIGN.CENTER)

    # Axis line
    R(sl, bx_start - Inches(0.05), CY + Inches(0.45),
      Inches(0.02), bar_max_h + Inches(0.45), MGRAY)
    R(sl, bx_start - Inches(0.05), bar_y_bottom,
      Inches(5.5), Inches(0.02), MGRAY)

    T(sl, px + Inches(0.12), bar_y_bottom + Inches(0.38), pw - Inches(0.24), Inches(0.28),
      "[ Illustrative values — run the app for real results ]",
      11, italic=True, color=STEEL, align=PP_ALIGN.CENTER)


def slide_forecast(prs):
    """Slide 13 — 7-Day Forecast Feature."""
    sl = new(prs)
    hf(sl, "7-Day Weather Forecast Feature", 13)

    # Left column
    lw = Inches(5.2)
    T(sl, ML, CY, lw, Inches(0.42), "HOW IT WORKS", 12, bold=True, color=BLUE)

    steps = [
        "Select station (Baguio / Manila)",
        "Tick the models to include",
        "Tick the variables to forecast",
        "Click Generate 7-Day Forecast",
        "View day-labeled table + multi-model line chart",
    ]
    numbered_list(sl, ML, CY + Inches(0.48), lw, Inches(2.5), steps,
                  size=14, color=DKTXT, num_color=BLUE)

    T(sl, ML, CY + Inches(3.1), lw, Inches(0.42),
      "TECHNICAL APPROACH", 12, bold=True, color=BLUE)
    tech = [
        "Loads scalers + model config from results cache",
        "Initialises last 45-day window from test set",
        "Iteratively predicts next day → appends to window",
        "ARIMA: .update(test_vals) then .predict(7)",
        "LR: same rolling-window inference",
        "Inverse-transforms all outputs to original units",
    ]
    bullet_list(sl, ML, CY + Inches(3.58), lw, Inches(2.3), tech,
                size=13, color=DKTXT, dot_color=NAVY)

    # Right: sample output mockup
    pw = CW - lw - Inches(0.3)
    px = ML + lw + Inches(0.3)

    # Table header
    R(sl, px, CY, pw, Inches(0.5), NAVY)
    T(sl, px + Inches(0.12), CY + Inches(0.08), pw - Inches(0.24), Inches(0.36),
      "Forecast — Precipitation (mm)  |  Baguio City", 13, bold=True, color=WHITE)

    # Table
    cols_ = ["Day", "Date", "GRU", "LSTM", "ARIMA"]
    col_ws = [Inches(1.1), Inches(1.5), Inches(0.9), Inches(0.9), Inches(0.9)]
    rows_ = [
        ["Monday",    "Jun 06", "1.2",  "0.9",  "1.5"],
        ["Tuesday",   "Jun 07", "4.8",  "3.2",  "5.1"],
        ["Wednesday", "Jun 08", "12.3", "10.7", "11.9"],
        ["Thursday",  "Jun 09", "2.1",  "1.8",  "2.4"],
        ["Friday",    "Jun 10", "0.3",  "0.4",  "0.6"],
        ["Saturday",  "Jun 11", "8.7",  "7.2",  "9.0"],
        ["Sunday",    "Jun 12", "5.4",  "4.9",  "6.1"],
    ]
    row_h = Inches(0.46)
    hdr_y = CY + Inches(0.5)
    cx_ = px
    for j, (ch_, cw_) in enumerate(zip(cols_, col_ws)):
        R(sl, cx_, hdr_y, cw_, row_h, BLUE)
        T(sl, cx_ + Inches(0.06), hdr_y + Inches(0.08), cw_ - Inches(0.12),
          row_h - Inches(0.16), ch_, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx_ += cw_

    for ri, row in enumerate(rows_):
        ry = hdr_y + row_h + ri * row_h
        cx_ = px
        for j, (val, cw_) in enumerate(zip(row, col_ws)):
            bg_c = LIGHT if ri % 2 == 0 else WHITE
            R(sl, cx_, ry, cw_, row_h, bg_c, line=MGRAY, lw=Pt(0.3))
            col_c = [C_GRU, C_LSTM, C_ARIMA][j - 2] if j >= 2 else DKTXT
            T(sl, cx_ + Inches(0.06), ry + Inches(0.08), cw_ - Inches(0.12),
              row_h - Inches(0.16), val, 12.5, color=col_c,
              align=PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT)
            cx_ += cw_

    T(sl, px, hdr_y + row_h * 8 + Inches(0.1), pw, Inches(0.34),
      "[ Sample values for illustration — actual forecasts generated from cached model data ]",
      11, italic=True, color=STEEL, align=PP_ALIGN.CENTER)


def slide_stack(prs):
    """Slide 14 — Technology Stack."""
    sl = new(prs)
    hf(sl, "Technology Stack", 14)

    stack = [
        ("Python 3.10+",            "Core language",                    NAVY),
        ("PyTorch 2.x",             "GRU · LSTM · SimpleRNN training",  C_GRU),
        ("pmdarima",                "auto_arima — AIC, stepwise",        C_ARIMA),
        ("scikit-learn",            "LinearRegression · MinMaxScaler",   C_RNN),
        ("Meteostat API",           "Historical daily weather data",     STAT_B),
        ("pandas · NumPy",          "Data wrangling & array ops",        STAT_G),
        ("Streamlit",               "Interactive web application UI",    STAT_P),
        ("Matplotlib",              "Prediction & metrics charts",       STAT_R),
        ("python-docx · pptx",      "Document & slide generation",       DARK),
        ("CUDA / MPS / CPU",        "Auto-detected GPU acceleration",    RGBColor(0x76, 0x44, 0x8a)),
    ]

    # Two columns
    mid = len(stack) // 2
    col_w = (CW - Inches(0.3)) / 2
    for col_idx, items in enumerate([stack[:mid], stack[mid:]]):
        cx = ML + col_idx * (col_w + Inches(0.3))
        for i, (tech, desc, col) in enumerate(items):
            ry = CY + i * Inches(0.6)
            R(sl, cx, ry, Inches(0.18), Inches(0.5), col)
            R(sl, cx + Inches(0.18), ry, col_w - Inches(0.18), Inches(0.5),
              LIGHT if i % 2 == 0 else OFFWH, line=MGRAY, lw=Pt(0.3))
            T(sl, cx + Inches(0.3), ry + Inches(0.06), Inches(2.4), Inches(0.38),
              tech, 14, bold=True, color=NAVY)
            T(sl, cx + Inches(2.78), ry + Inches(0.06),
              col_w - Inches(2.82), Inches(0.38), desc, 13.5, color=DKTXT)

    # Deployment note
    note_y = CY + mid * Inches(0.6) + Inches(0.18)
    R(sl, ML, note_y, CW, Inches(0.82), NAVY)
    T(sl, ML + Inches(0.3), note_y + Inches(0.12), CW - Inches(0.6), Inches(0.6),
      "Deployed on  Streamlit Community Cloud  ·  "
      "thesis-gru-evaluation-in-weather-forecasting.streamlit.app  ·  "
      "Source: github.com/EzekelCastro/Thesis-GRU-Evaluation-in-Weather-Forecasting",
      13, color=STEEL)


def slide_conclusion(prs):
    """Slide 15 — Conclusion & Future Work."""
    sl = new(prs)
    hf(sl, "Conclusion & Future Work", 15)

    lw = Inches(6.0)

    T(sl, ML, CY, lw, Inches(0.42), "KEY CONTRIBUTIONS", 12, bold=True, color=BLUE)
    contribs = [
        "End-to-end benchmarking of 5 ML models on Philippine weather data",
        "Station-level comparison: Baguio City (highland) vs Manila (lowland)",
        "Variable Importance analysis: leave-one-out ablation across all targets × both stations",
        "Dynamic variable selection from all 10 Meteostat inputs for flexible evaluation",
        "Interactive web system with Variable Importance, forecasting, and cached results",
    ]
    bullet_list(sl, ML, CY + Inches(0.48), lw, Inches(2.5), contribs,
                size=14, color=DKTXT, dot_color=BLUE)

    T(sl, ML, CY + Inches(3.1), lw, Inches(0.42),
      "EXPECTED OUTCOMES", 12, bold=True, color=BLUE)
    outcomes = [
        "GRU expected to outperform SimpleRNN and ARIMA on most variables",
        "LSTM competitive — architecture comparison between GRU and LSTM is central",
        "Pressure and Temperature: highest model accuracy across all architectures",
        "Precipitation remains the most challenging variable to forecast accurately",
    ]
    bullet_list(sl, ML, CY + Inches(3.58), lw, Inches(2.1), outcomes,
                size=14, color=DKTXT, dot_color=NAVY)

    # Right column: future work
    rw = CW - lw - Inches(0.3)
    rx = ML + lw + Inches(0.3)
    R(sl, rx, CY, rw, Inches(5.6), NAVY)
    T(sl, rx + Inches(0.2), CY + Inches(0.18), rw - Inches(0.4), Inches(0.38),
      "FUTURE WORK", 12, bold=True, color=STEEL)

    future = [
        "Incorporate additional stations\n(Cebu, Davao, Zamboanga)",
        "Add Transformer-based model\n(Temporal Fusion Transformer)",
        "Seasonal decomposition as\npreprocessing step for prcp",
        "Ensemble prediction combining\nmultiple model outputs",
        "Mobile-friendly app with\npush notification forecasts",
    ]
    for i, item in enumerate(future):
        fy = CY + Inches(0.65) + i * Inches(0.98)
        block(sl, rx + Inches(0.2), fy, rw - Inches(0.4), Inches(0.86),
              BLUE if i % 2 == 0 else RGBColor(0x1a, 0x52, 0x76),
              [(item, 13, False, WHITE)], v_anchor=MSO_ANCHOR.MIDDLE)


def slide_thankyou(prs):
    """Slide 16 — Thank You / Q&A."""
    sl = new(prs)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = DARK

    T(sl, ML, Inches(1.3), CW, Inches(0.52),
      "■  THANK YOU  ■", 14, bold=True, color=STEEL, align=PP_ALIGN.CENTER)

    T(sl, ML, Inches(1.95), CW, Inches(0.92),
      "Questions & Discussion", 44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    R(sl, ML + Inches(2.5), Inches(3.0), CW - Inches(5), Inches(0.05), BLUE)

    # Authors
    T(sl, ML, Inches(3.2), CW, Inches(0.5),
      "Castro, E. M.  ·  Millan, D. J. S.  ·  Mondoñedo, J. E.",
      17, color=RGBColor(0xad, 0xc8, 0xe0), align=PP_ALIGN.CENTER)

    T(sl, ML, Inches(3.76), CW, Inches(0.4),
      "University of the Cordilleras  ·  College of CITCS  ·  B.S. Computer Science",
      14, color=STEEL, align=PP_ALIGN.CENTER)

    # Links
    R(sl, ML + Inches(1.5), Inches(4.4), CW - Inches(3.0), Inches(1.2), NAVY)
    T(sl, ML + Inches(1.7), Inches(4.58), CW - Inches(3.4), Inches(0.4),
      "🔗  thesis-gru-evaluation-in-weather-forecasting.streamlit.app",
      14, color=RGBColor(0x5d, 0xad, 0xe8), align=PP_ALIGN.CENTER)
    T(sl, ML + Inches(1.7), Inches(5.02), CW - Inches(3.4), Inches(0.4),
      "💻  github.com/EzekelCastro/Thesis-GRU-Evaluation-in-Weather-Forecasting",
      14, color=RGBColor(0x5d, 0xad, 0xe8), align=PP_ALIGN.CENTER)

    # Model colour strip
    models = [("GRU", C_GRU), ("LSTM", C_LSTM), ("SimpleRNN", C_RNN),
              ("Linear Regression", C_LR), ("ARIMA", C_ARIMA)]
    strip_w = CW / len(models)
    for i, (name, col) in enumerate(models):
        chip(sl, ML + i * strip_w, Inches(6.0), strip_w, Inches(0.36),
             name, col, size=13)

    R(sl, 0, SH - FTR, SW, FTR, RGBColor(0x17, 0x20, 0x2a))
    T(sl, ML, SH - FTR + Inches(0.05), SW - ML * 2, Inches(0.28),
      f"University of the Cordilleras  ·  Baguio City  ·  {YEAR}",
      10, color=STEEL, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs)
    slide_glance(prs)
    slide_problem(prs)
    slide_data(prs)
    slide_variables(prs)
    slide_models(prs)
    slide_architecture(prs)
    slide_preprocessing(prs)
    slide_ui_overview(prs)
    slide_metrics(prs)
    slide_results_predictions(prs)
    slide_results_metrics(prs)
    slide_forecast(prs)
    slide_stack(prs)
    slide_conclusion(prs)
    slide_thankyou(prs)

    prs.save("thesis_presentation.pptx")
    print(f"Saved: thesis_presentation.pptx  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
